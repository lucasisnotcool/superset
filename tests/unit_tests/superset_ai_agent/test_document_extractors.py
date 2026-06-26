# Licensed to the Apache Software Foundation (ASF) under one
# or more contributor license agreements.  See the NOTICE file
# distributed with this work for additional information
# regarding copyright ownership.  The ASF licenses this file
# to you under the Apache License, Version 2.0 (the
# "License"); you may not use this file except in compliance
# with the License.  You may obtain a copy of the License at
#
#   http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing,
# software distributed under the License is distributed on an
# "AS IS" BASIS, WITHOUT WARRANTIES OR CONDITIONS OF ANY
# KIND, either express or implied.  See the License for the
# specific language governing permissions and limitations
# under the License.

from __future__ import annotations

import importlib.util
import io

import pytest

from superset_ai_agent.semantic_layer.extractors import (
    _DOCX_CONTENT_TYPE,
    _PPTX_CONTENT_TYPE,
    _rows_to_markdown_table,
    _XLSX_CONTENT_TYPE,
    CompositeDocumentExtractor,
    NeedsOcrError,
)

_PDF_AVAILABLE = importlib.util.find_spec("pypdf") is not None
_DOCX_AVAILABLE = importlib.util.find_spec("docx") is not None
_XLSX_AVAILABLE = importlib.util.find_spec("openpyxl") is not None
_PPTX_AVAILABLE = importlib.util.find_spec("pptx") is not None


def _extract(content_type: str, content: bytes) -> str:
    return CompositeDocumentExtractor().extract_text(
        filename="doc",
        content_type=content_type,
        content=content,
    )


def test_html_extraction_strips_tags_and_scripts() -> None:
    html = (
        b"<html><head><title>x</title><style>p{}</style></head>"
        b"<body><h1>Revenue</h1><p>Orders join customers.</p>"
        b"<script>evil()</script></body></html>"
    )
    text = _extract("text/html", html)
    assert "Revenue" in text
    assert "Orders join customers." in text
    assert "evil" not in text  # script content dropped
    assert "{}" not in text  # style content dropped


def test_html_content_type_parameters_are_normalized() -> None:
    text = _extract("text/html; charset=utf-8", b"<p>hello</p>")
    assert text == "hello"


def test_unsupported_type_still_raises() -> None:
    with pytest.raises(ValueError, match="Unsupported document content type"):
        _extract("image/png", b"\x89PNG")


def test_pdf_missing_dependency_raises_clear_error() -> None:
    if _PDF_AVAILABLE:
        pytest.skip("pypdf installed; missing-dependency path not exercised")
    with pytest.raises(ValueError, match="pypdf"):
        _extract("application/pdf", b"%PDF-1.4 ...")


def test_docx_missing_dependency_raises_clear_error() -> None:
    if _DOCX_AVAILABLE:
        pytest.skip("python-docx installed; missing-dependency path not exercised")
    with pytest.raises(ValueError, match="python-docx"):
        _extract(_DOCX_CONTENT_TYPE, b"PK\x03\x04 ...")


@pytest.mark.skipif(not _DOCX_AVAILABLE, reason="python-docx not installed")
def test_docx_round_trip_extracts_paragraphs() -> None:
    import docx

    document = docx.Document()
    document.add_paragraph("First line.")
    document.add_paragraph("Second line.")
    buffer = io.BytesIO()
    document.save(buffer)

    text = _extract(_DOCX_CONTENT_TYPE, buffer.getvalue())
    assert "First line." in text
    assert "Second line." in text


@pytest.mark.skipif(
    not (_PDF_AVAILABLE and importlib.util.find_spec("reportlab")),
    reason="pypdf and reportlab required to build and read a text PDF",
)
def test_pdf_round_trip_extracts_text() -> None:
    from reportlab.pdfgen import canvas

    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer)
    pdf.drawString(72, 720, "Hello from PDF")
    pdf.save()

    text = _extract("application/pdf", buffer.getvalue())
    assert "Hello from PDF" in text


# --- Markdown-table helper (Step 2) ---------------------------------------


def test_rows_to_markdown_table_renders_header_and_separator() -> None:
    table = _rows_to_markdown_table([["a", "b"], ["1", "2"]])
    assert table.splitlines() == [
        "| a | b |",
        "| --- | --- |",
        "| 1 | 2 |",
    ]


def test_rows_to_markdown_table_pads_ragged_and_escapes_pipes() -> None:
    table = _rows_to_markdown_table([["h1", "h2"], ["a|b"]])
    lines = table.splitlines()
    assert lines[0] == "| h1 | h2 |"
    # ragged row padded to width 2, pipe escaped
    assert lines[2] == "| a\\|b |  |"


def test_rows_to_markdown_table_empty_returns_blank() -> None:
    assert _rows_to_markdown_table([]) == ""
    assert _rows_to_markdown_table([["", "  "]]) == ""


# --- CSV upgrade (Step 2) -------------------------------------------------


def test_csv_renders_markdown_table() -> None:
    text = _extract("text/csv", b"name,metric\nrevenue,sum(amount)\n")
    assert text.splitlines() == [
        "| name | metric |",
        "| --- | --- |",
        "| revenue | sum(amount) |",
    ]


def test_csv_empty_returns_blank() -> None:
    assert _extract("text/csv", b"") == ""


# --- XLSX (Step 3) --------------------------------------------------------


def test_xlsx_missing_dependency_raises_clear_error() -> None:
    if _XLSX_AVAILABLE:
        pytest.skip("openpyxl installed; missing-dependency path not exercised")
    with pytest.raises(ValueError, match="openpyxl"):
        _extract(_XLSX_CONTENT_TYPE, b"PK\x03\x04 ...")


@pytest.mark.skipif(not _XLSX_AVAILABLE, reason="openpyxl not installed")
def test_xlsx_round_trip_renders_per_sheet_tables() -> None:
    import openpyxl

    workbook = openpyxl.Workbook()
    first = workbook.active
    first.title = "Metrics"
    first.append(["name", "definition"])
    first.append(["revenue", "sum of amount"])
    second = workbook.create_sheet("Dims")
    second.append(["dim", "grain"])
    second.append(["date", "day"])
    buffer = io.BytesIO()
    workbook.save(buffer)

    text = _extract(_XLSX_CONTENT_TYPE, buffer.getvalue())
    assert "## Sheet: Metrics" in text
    assert "## Sheet: Dims" in text
    assert "| name | definition |" in text
    assert "| revenue | sum of amount |" in text
    # sheets are separated into distinct blocks
    assert text.index("## Sheet: Metrics") < text.index("## Sheet: Dims")


# --- PPTX (Step 4) --------------------------------------------------------


def test_pptx_missing_dependency_raises_clear_error() -> None:
    if _PPTX_AVAILABLE:
        pytest.skip("python-pptx installed; missing-dependency path not exercised")
    with pytest.raises(ValueError, match="python-pptx"):
        _extract(_PPTX_CONTENT_TYPE, b"PK\x03\x04 ...")


@pytest.mark.skipif(not _PPTX_AVAILABLE, reason="python-pptx not installed")
def test_pptx_round_trip_extracts_slides_text_and_tables() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    blank = presentation.slide_layouts[6]

    slide_one = presentation.slides.add_slide(blank)
    box = slide_one.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Quarterly KPIs"

    slide_two = presentation.slides.add_slide(blank)
    graphic = slide_two.shapes.add_table(
        2, 2, Inches(1), Inches(1), Inches(4), Inches(2)
    )
    table = graphic.table
    table.cell(0, 0).text = "kpi"
    table.cell(0, 1).text = "target"
    table.cell(1, 0).text = "churn"
    table.cell(1, 1).text = "5%"

    buffer = io.BytesIO()
    presentation.save(buffer)

    text = _extract(_PPTX_CONTENT_TYPE, buffer.getvalue())
    assert "## Slide 1" in text
    assert "Quarterly KPIs" in text
    assert "## Slide 2" in text
    assert "| kpi | target |" in text
    assert "| churn | 5% |" in text


# --- PDF image-only detection (Step 5) ------------------------------------


@pytest.mark.skipif(not _PDF_AVAILABLE, reason="pypdf not installed")
def test_pdf_without_text_layer_raises_needs_ocr(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import superset_ai_agent.semantic_layer.extractors as extractors

    class _Page:
        def extract_text(self) -> str:
            return "   "

    class _Reader:
        def __init__(self, *_args: object, **_kwargs: object) -> None:
            self.pages = [_Page(), _Page()]

    monkeypatch.setattr("pypdf.PdfReader", _Reader)
    with pytest.raises(NeedsOcrError):
        extractors._extract_pdf(b"%PDF-1.4 image only")
