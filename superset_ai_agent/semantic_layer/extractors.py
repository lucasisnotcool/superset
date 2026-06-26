# Licensed to the Apache Software Foundation (ASF) under one
# or more contributor license agreements.  See the NOTICE file
# distributed with this work for additional information
# regarding copyright ownership.  The ASF licenses this file
# to you under the Apache License, Version 2.0 (the
# "License"); you may not use this file except in compliance
# with the License.  You may obtain a copy of the License at
#
#   http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing,
# software distributed under the License is distributed on an
# "AS IS" BASIS, WITHOUT WARRANTIES OR CONDITIONS OF ANY
# KIND, either express or implied.  See the License for the
# specific language governing permissions and limitations
# under the License.

from __future__ import annotations

import csv
import io
import json  # noqa: TID251 - standalone agent: stdlib JSON for document parsing
from html.parser import HTMLParser
from typing import Protocol

#: Office Open XML Word document MIME type (.docx).
_DOCX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
)
#: Office Open XML Spreadsheet MIME type (.xlsx).
_XLSX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
#: Office Open XML Presentation MIME type (.pptx).
_PPTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


class DocumentExtractionError(ValueError):
    """Base for extraction failures.

    Subclasses ``ValueError`` so the existing endpoint 400-handling and
    ``create_document``'s ``except`` paths keep working unchanged.
    """


class NeedsOcrError(DocumentExtractionError):
    """Raised when a file has no extractable text layer and needs OCR.

    OCR itself is out of scope; raising this is the seam. ``create_document`` maps
    it to ``status="needs_ocr"`` (not ``"error"``) so a future OCR backend can pick
    these documents up. The original bytes are already stored, so OCR can run later
    without a re-upload.
    """


class DocumentExtractor(Protocol):
    """Extract safe plain text from supported document types."""

    def extract_text(self, *, filename: str, content_type: str, content: bytes) -> str:
        """Return plain text extracted from uploaded content."""


class CompositeDocumentExtractor:
    """Extractor for text, Markdown, CSV, JSON, HTML, PDF, DOCX, XLSX, and PPTX.

    Tabular content (CSV, XLSX, PPTX tables) is rendered as GitHub-flavoured
    Markdown tables, and workbook/deck provenance is encoded as ``## Sheet:`` /
    ``## Slide n`` section headers so the downstream chunker keeps each sheet/slide
    as a retrievable unit (canonical-text contract; see
    ``document_format_tier1_plan.md`` A.1).

    HTML/CSV/JSON extraction is dependency-free (stdlib). PDF, DOCX, XLSX and PPTX
    rely on optional packages (``pypdf`` / ``python-docx`` / ``openpyxl`` /
    ``python-pptx``); when a package is missing the extractor raises ``ValueError``
    so ``create_document`` records the document with ``status="error"`` and a clear
    message, rather than crashing the upload (governance: degrade closed). An
    image-only PDF (no text layer) raises ``NeedsOcrError`` so it is tagged
    ``needs_ocr`` rather than stored with empty text.
    """

    def extract_text(self, *, filename: str, content_type: str, content: bytes) -> str:
        normalized_type = normalize_content_type(content_type)
        if normalized_type in {"text/plain", "text/markdown"}:
            return _decode_text(content)
        if normalized_type == "text/csv":
            return _extract_csv(content)
        if normalized_type == "application/json":
            return _extract_json(content)
        if normalized_type == "text/html":
            return _extract_html(content)
        if normalized_type == "application/pdf":
            return _extract_pdf(content)
        if normalized_type == _DOCX_CONTENT_TYPE:
            return _extract_docx(content)
        if normalized_type == _XLSX_CONTENT_TYPE:
            return _extract_xlsx(content)
        if normalized_type == _PPTX_CONTENT_TYPE:
            return _extract_pptx(content)
        raise ValueError(f"Unsupported document content type: {content_type}")


def normalize_content_type(content_type: str) -> str:
    """Return a comparable content type without parameters."""

    return content_type.split(";", 1)[0].strip().lower()


def _decode_text(content: bytes) -> str:
    text = content.decode("utf-8", errors="replace")
    return _strip_nul(text)


def _extract_csv(content: bytes) -> str:
    text = _decode_text(content)
    reader = csv.reader(io.StringIO(text))
    rows = [list(row) for row in reader]
    return _rows_to_markdown_table(rows)


def _rows_to_markdown_table(rows: list[list[str]]) -> str:
    """Render rows as a GitHub-flavoured Markdown table (first row = header).

    Drops fully-empty rows, pads ragged rows to a rectangular grid, and escapes
    pipes/newlines in cells. Returns ``""`` when there is no data. A single row is
    still emitted as a header-only table.
    """

    cleaned = [row for row in rows if any(str(cell).strip() for cell in row)]
    if not cleaned:
        return ""
    width = max(len(row) for row in cleaned)
    grid = [
        [_escape_cell(row[i]) if i < len(row) else "" for i in range(width)]
        for row in cleaned
    ]
    header = grid[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in range(width)) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in grid[1:])
    return "\n".join(lines)


def _escape_cell(cell: object) -> str:
    return str(cell).replace("|", "\\|").replace("\n", " ").replace("\r", " ").strip()


def _extract_json(content: bytes) -> str:
    parsed = json.loads(_decode_text(content))
    return json.dumps(parsed, indent=2, sort_keys=True)


def _strip_nul(text: str) -> str:
    return text.replace("\x00", "")


class _HtmlTextExtractor(HTMLParser):
    """Collect visible text from HTML, dropping script/style and block-spacing tags."""

    _SKIP = {"script", "style", "head", "title", "meta"}
    _BLOCK = {
        "p",
        "br",
        "div",
        "li",
        "tr",
        "h1",
        "h2",
        "h3",
        "h4",
        "h5",
        "h6",
        "section",
        "article",
    }

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._parts: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: object) -> None:
        if tag in self._SKIP:
            self._skip_depth += 1
        elif tag in self._BLOCK:
            self._parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP and self._skip_depth:
            self._skip_depth -= 1
        elif tag in self._BLOCK:
            self._parts.append("\n")

    def handle_data(self, data: str) -> None:
        if self._skip_depth == 0 and data.strip():
            self._parts.append(data)

    def get_text(self) -> str:
        lines = [line.strip() for line in "".join(self._parts).splitlines()]
        return "\n".join(line for line in lines if line)


def _extract_html(content: bytes) -> str:
    parser = _HtmlTextExtractor()
    parser.feed(_decode_text(content))
    parser.close()
    return parser.get_text()


def _extract_pdf(content: bytes) -> str:
    try:
        from pypdf import PdfReader  # optional dependency (ignore_missing_imports)
    except ImportError as ex:
        raise ValueError(
            "PDF extraction requires the 'pypdf' package to be installed."
        ) from ex
    reader = PdfReader(io.BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    non_empty = [page.strip() for page in pages if page.strip()]
    if len(reader.pages) > 0 and not non_empty:
        # The PDF has pages but no extractable text layer (scanned / image-only).
        # Tag for OCR instead of saving an empty document. OCR is out of scope here.
        raise NeedsOcrError("PDF has no extractable text layer; OCR required.")
    return _strip_nul("\n\n".join(non_empty))


def _extract_docx(content: bytes) -> str:
    try:
        import docx  # optional dependency: python-docx (ignore_missing_imports)
    except ImportError as ex:
        raise ValueError(
            "DOCX extraction requires the 'python-docx' package to be installed."
        ) from ex
    document = docx.Document(io.BytesIO(content))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs]
    return _strip_nul("\n".join(text for text in paragraphs if text))


def _extract_xlsx(content: bytes) -> str:
    try:
        import openpyxl  # optional dependency (ignore_missing_imports)
    except ImportError as ex:
        raise ValueError(
            "XLSX extraction requires the 'openpyxl' package to be installed."
        ) from ex
    # read_only streams large sheets; data_only yields last-computed values
    # (formulas themselves are not preserved, which is acceptable for text recall).
    workbook = openpyxl.load_workbook(
        io.BytesIO(content), read_only=True, data_only=True
    )
    try:
        blocks: list[str] = []
        for worksheet in workbook.worksheets:
            rows = [
                ["" if cell is None else str(cell) for cell in row]
                for row in worksheet.iter_rows(values_only=True)
            ]
            table = _rows_to_markdown_table(rows)
            if table:
                blocks.append(f"## Sheet: {worksheet.title}\n\n{table}")
    finally:
        workbook.close()
    return _strip_nul("\n\n".join(blocks))


def _extract_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation  # optional dependency: python-pptx
    except ImportError as ex:
        raise ValueError(
            "PPTX extraction requires the 'python-pptx' package to be installed."
        ) from ex
    presentation = Presentation(io.BytesIO(content))
    blocks: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                table = _rows_to_markdown_table(rows)
                if table:
                    parts.append(table)
            elif shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
        if parts:
            blocks.append(f"## Slide {index}\n\n" + "\n\n".join(parts))
    return _strip_nul("\n\n".join(blocks))
